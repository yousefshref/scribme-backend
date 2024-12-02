import base64
import datetime
import tempfile
from django.conf import settings
from PIL import Image, ImageEnhance
from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework import status
from pptx import Presentation
from fitz import open as open_pdf
from .serializers import ImageDescriptionSerializer
from pathlib import Path
import os
import io


# Global image description count and last recharge time
image_description_count = 20
last_recharge_time = datetime.datetime.now()

from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework import status
from PIL import Image
import pytesseract

import requests

from rest_framework.decorators import api_view
import fitz  # PyMuPDF for working with PDFs


from . import models, serializers


pytesseract.pytesseract.tesseract_cmd = 'C:\\Program Files\\Tesseract-OCR\\tesseract.exe'


def describe_image_with_gpt(base64_image, prompt_text="Describe this image"):
    # Set up the payload for OpenAI API
    api_key = "sk-proj-8kPWAezn8UH2b7muCR0kT3BlbkFJBZTrz86pU5tzXnbw9zGy"  # Make sure to set this in your settings
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {api_key}",
    }

    payload = {
        "model": "gpt-4o-2024-08-06",
        "messages": [
            {
            "role": "user",
            "content": [
                {
                "type": "text",
                "text": prompt_text 
                },
                {
                "type": "image_url",
                "image_url": {
                    "url": f"data:image/jpeg;base64,{base64_image}",
                    
                }
                }
            ]
            }
        ],
        "max_tokens": 325
        }

    # Send request to OpenAI API
    response = requests.post(
        "https://api.openai.com/v1/chat/completions",
        headers=headers,
        json=payload,
    )
    response.raise_for_status()  # Check if the request was successful
    response_json = response.json()

    # Extract the description content from the response
    description = response_json["choices"][0]["message"]["content"]

    return description





class DescribeImageView(APIView):
    def post(self, request):
        image_file = request.FILES.get("image")  # Retrieve the uploaded file
        language = request.data.get("language", "English")

        if not image_file:
            return Response({"error": "Image file is required."}, status=status.HTTP_400_BAD_REQUEST)

        try:
            # Open the image file using PIL
            image = Image.open(image_file)

            # Set OCR language configuration based on requested language
            custom_config = '--psm 1 --oem 1'
            if language == "Arabic":
                custom_config += ' -l ara'
            elif language == "Spanish":
                custom_config += ' -l spa'
            else:
                custom_config += ' -l eng'

            # Perform OCR on the image
            # text = pytesseract.image_to_string(image, config=custom_config)
            # Convert the image file to base64
            buffered = io.BytesIO()
            image.save(buffered, format="JPEG")
            base64_image = base64.b64encode(buffered.getvalue()).decode("utf-8")

            # Define prompt text based on the language
            if language == "English":
                prompt_text = "Describe this image in detail."
            elif language == "Arabic":
                prompt_text = " ."
            elif language == "Spanish":
                prompt_text = "Describe esta imagen en detalle."
            else:
                prompt_text = "Describe this image in detail."

            # Call the function to describe the image with GPT
            text = describe_image_with_gpt(base64_image, prompt_text)

            return Response({"description": text}, status=status.HTTP_200_OK)

        except Exception as e:
            return Response({"error": str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


class ExtractTextFromPDFView(APIView):
    def post(self, request):
        pdf_file = request.FILES.get("pdf_file")
        ocr_option = request.data.get("ocr", False)  # If OCR is requested
        image_description_option = request.data.get("image_description", False)  # If image description is requested
        
        if not pdf_file:
            return Response({"error": "PDF file is required."}, status=status.HTTP_400_BAD_REQUEST)

        # Save the PDF to a temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
            temp_pdf.write(pdf_file.read())
            temp_pdf_path = temp_pdf.name

        try:
            text_content = ""
            with open_pdf(temp_pdf_path) as pdf_document:
                for page_number, page in enumerate(pdf_document, start=1):
                    text_content += f"Page {page_number}:\n"
                    text_content += page.get_text("text") + "\n"

                    # If OCR or image description is enabled, process images on the page
                    if ocr_option or image_description_option:
                        images = page.get_images(full=True)
                        for img_index, img in enumerate(images):
                            xref = img[0]
                            base_image = pdf_document.extract_image(xref)
                            image_bytes = base_image["image"]

                            # Save extracted image to a temporary file
                            with tempfile.NamedTemporaryFile(delete=False, suffix=".jpg") as temp_image:
                                temp_image.write(image_bytes)
                                temp_image_path = temp_image.name

                            # Perform OCR on image if OCR option is enabled
                            if ocr_option:
                                ocr_text = self.perform_ocr(temp_image_path)
                                text_content += f"\n OCR Text from image on page {page_number}: {ocr_text}\n"

                            # Describe image with GPT if image description option is enabled
                            if image_description_option:
                                gpt_description = describe_image_with_gpt(
                                    base64_image=base64.b64encode(image_bytes).decode("utf-8"),
                                    prompt_text="Describe this image in detail.",
                                )
                                text_content += f"\n Image description on page {page_number}: {gpt_description}\n"

                            # Clean up the temporary image file
                            os.remove(temp_image_path)

            os.remove(temp_pdf_path)
            return Response({"text_content": text_content}, status=status.HTTP_200_OK)

        except Exception as e:
            # Ensure the temporary file is deleted in case of error
            if os.path.exists(temp_pdf_path):
                os.remove(temp_pdf_path)
            return Response({"error": str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

    def perform_ocr(self, image_path):
        # Helper function to perform OCR on an image
        image = Image.open(image_path)
        ocr_text = pytesseract.image_to_string(image, lang='eng')
        return ocr_text








from pptx import Presentation
from rest_framework.parsers import MultiPartParser
from rest_framework.views import APIView
from rest_framework.response import Response
import base64
import comtypes.client
import tempfile

def convert_ppt_to_pptx(ppt_file_path):
    """
    Convert a .ppt file to .pptx using PowerPoint application on Windows.
    
    Args:
        ppt_file_path (str): The path to the .ppt file.

    Returns:
        str: The path to the converted .pptx file.
    """
    try:
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.Visible = 1  # Make PowerPoint visible (optional)
        
        # Open the .ppt file
        presentation = powerpoint.Presentations.Open(ppt_file_path)
        
        # Generate a temporary path for the converted .pptx
        pptx_file_path = ppt_file_path + 'x'  # Convert .ppt to .pptx
        
        # Save as .pptx
        presentation.SaveAs(pptx_file_path, 24)  # 24 is for .pptx format
        presentation.Close()
        powerpoint.Quit()
        
        return pptx_file_path

    except Exception as e:
        raise Exception(f"Error converting .ppt to .pptx: {e}")


def extract_content_from_pptx(presentation):
    """
    Extracts text and images from each slide in a PowerPoint presentation.

    Args:
        presentation (Presentation): The PowerPoint presentation object.

    Returns:
        list: A list of dictionaries, each representing a slide with its text and images.
    """
    slides_content = []

    try:
        # Iterate through slides
        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_data = {"slide_number": slide_index, "texts": [], "images": []}

            # Extract text from shapes
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        slide_data["texts"].append(paragraph.text.strip())

                # Extract images
                if hasattr(shape, "image"):
                    image_stream = shape.image.blob
                    image_base64 = base64.b64encode(image_stream).decode("utf-8")
                    slide_data["images"].append(image_base64)

            slides_content.append(slide_data)

    except Exception as e:
        print(f"Error extracting content: {e}")

    return slides_content



class PptxProcessorAPIView(APIView):
    parser_classes = [MultiPartParser]

    def post(self, request, *args, **kwargs):
        pptx_file = request.FILES.get("file")
        if not pptx_file:
            return Response({"error": "No file uploaded."})

        try:
            # Save the file to a temporary location
            with tempfile.NamedTemporaryFile(delete=False) as temp_file:
                temp_file.write(pptx_file.read())
                temp_file_path = temp_file.name

            # Check if the file is a .ppt or .pptx
            if temp_file_path.lower().endswith('.ppt'):
                # Convert .ppt to .pptx
                pptx_file_path = convert_ppt_to_pptx(temp_file_path)
            else:
                pptx_file_path = temp_file_path

            # Load the PowerPoint presentation
            presentation = Presentation(pptx_file_path)

            # Extract all text and images from the presentation
            slides_content = extract_content_from_pptx(presentation)

            # Describe images using GPT and replace Base64 strings with descriptions
            for slide in slides_content:
                described_images = []
                for image_base64 in slide["images"]:
                    description = describe_image_with_gpt(image_base64)
                    described_images.append(description)
                slide["images"] = described_images

            return Response({"slides": slides_content})

        except Exception as e:
            return Response({"error": str(e)})


# class PptxProcessorAPIView(APIView):
#     parser_classes = [MultiPartParser]

#     def post(self, request, *args, **kwargs):
#         pptx_file = request.FILES.get("file")
#         if not pptx_file:
#             return Response({"error": "No file uploaded."})

#         try:
#             # Load the PowerPoint presentation
#             presentation = Presentation(pptx_file)

#             # Extract all text and images from the presentation
#             slides_content = extract_content_from_pptx(presentation)

#             # Describe images using GPT and replace Base64 strings with descriptions
#             for slide in slides_content:
#                 described_images = []
#                 for image_base64 in slide["images"]:
#                     description = describe_image_with_gpt(image_base64)
#                     described_images.append(description)
#                 slide["images"] = described_images

#             return Response({"slides": slides_content})

#         except Exception as e:
#             return Response({"error": str(e)})




# CRUD history
@api_view(["POST"])
def create_history(request):
    serializer = serializers.HistorySerializer(data=request.data)
    if serializer.is_valid():
        serializer.save()
        return Response(serializer.data, status=status.HTTP_201_CREATED)
    return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

@api_view(["GET"])
def get_history(request, user_id):
    history = models.History.objects.filter(user=user_id)
    serializer = serializers.HistorySerializer(history, many=True)
    return Response(serializer.data)

@api_view(["GET"])
def get_history_by_id(request, pk):
    history = models.History.objects.get(pk=pk)
    serializer = serializers.HistorySerializer(history)
    return Response(serializer.data)


