import wx
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import requests
import pytesseract
from PIL import Image
from PIL import ImageEnhance
import fitz
import io
import sys
import os
from pathlib import Path
import tempfile
import base64
import requests
import threading
import time
import keyboard
import datetime
import json
import config

# global variables
image_description_count = 20
last_recharge_time = datetime.datetime.now()

pytesseract.pytesseract.tesseract_cmd = 'tesseract.exe'                

APP_DATA_DIR = Path(os.getenv('APPDATA')) / 'ScribeMe'

PERSISTENCE_FILE = APP_DATA_DIR / 'image_description_data.json'

APP_DATA_DIR.mkdir(parents=True, exist_ok=True)

def save_image_description_data():
    global image_description_count, last_recharge_time

    data = {
        'image_description_count': image_description_count,
        'last_recharge_time': last_recharge_time.isoformat()  # Convert datetime to string
    }

    # Write data to the file
    try:
        with open(PERSISTENCE_FILE, 'w') as f:
            json.dump(data, f, indent=4)
        print(f"Data successfully saved to {PERSISTENCE_FILE}")
    except Exception as e:
        print(f"Error saving data to {PERSISTENCE_FILE}: {e}")

def load_image_description_data():
    global image_description_count, last_recharge_time
    if os.path.exists(PERSISTENCE_FILE):
        try:
            with open(PERSISTENCE_FILE, 'r') as f:
                data = json.load(f)
                image_description_count = data.get('image_description_count', 20)
                last_recharge_time_str = data.get('last_recharge_time', None)
                last_recharge_time = datetime.datetime.fromisoformat(last_recharge_time_str) if last_recharge_time_str else datetime.datetime.now()
                now = datetime.datetime.now()
                elapsed_time = (now - last_recharge_time).total_seconds() / 3600.0  # convert to hours
                if elapsed_time >= 4:
                    # Calculate how many recharges have occurred
                    recharge_count = int(elapsed_time // 4)
                    image_description_count = min(20, image_description_count + (recharge_count * 5))
                    last_recharge_time = last_recharge_time + datetime.timedelta(hours=recharge_count * 4)
                print(data)
            print(f"Data successfully loaded from {PERSISTENCE_FILE}")
        except Exception as e:
            print(f"Error loading data from {PERSISTENCE_FILE}: {e}")
            image_description_count = 20
            last_recharge_time = datetime.datetime.now()
    else:
        print(f"{PERSISTENCE_FILE} does not exist. Initializing with default values.")
        image_description_count = 20
        last_recharge_time = datetime.datetime.now()

def recharge_image_descriptions():
    global image_description_count, last_recharge_time
    while True:
        now = datetime.datetime.now()
        elapsed_time = (now - last_recharge_time).total_seconds() / 3600.0  # convert to hours
        print(elapsed_time)
        if elapsed_time >= 4:
            image_description_count += 5
            if image_description_count > 20:
                image_description_count = 20  # cap at 20
            last_recharge_time = now
            save_image_description_data()  # Save the data after recharging
        
        time.sleep(3600)  # check every hour

# Start the recharge thread
threading.Thread(target=recharge_image_descriptions, daemon=True).start()

    
def describe_image(image_path, language="English"):
    # API key
    api_key = config.api_key
    
    # Encode the image
    with open(image_path, "rb") as image_file:
        base64_image = base64.b64encode(image_file.read()).decode('utf-8')

    # Define headers
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {api_key}"
    }
    if language == "English":
        prompt_text = "What's in this image?"
    elif language == "Arabic":
        prompt_text = "ما هو في هذه الصورة؟"
    elif language == "Spanish":
        prompt_text = "¿Qué hay en esta imagen?"
    
    

    # Define payload
    
    payload = {
        "model": "gpt-4o-2024-08-06",
        "messages": [
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": prompt_text  
                    },
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:image/jpeg;base64,{base64_image}",
                            
                        }
                    }
                ]
            }
        ],
        "max_tokens": 325
    }

    # Send request to OpenAI API
    response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload)
    response_json = response.json()

    # Extract the description content
    description = response_json['choices'][0]['message']['content']
    
    # Return only the description
    return description

def take_screenshot():
    import pyautogui

    with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_file:
        screenshot_path = temp_file.name
        print(screenshot_path)
        
    try:
        # Take the screenshot using pyautogui
        screenshot = pyautogui.screenshot()
        # Save the screenshot to the temporary file
        screenshot.save(screenshot_path)
        return screenshot_path
    except Exception as e:
        print(f"An error occurred: {e}")

def preprocess_image(image_path):
    image = Image.open(image_path)
    # Convert to grayscale
    image = image.convert('L')
    # Enhance contrast
    enhancer = ImageEnhance.Contrast(image)
    image = enhancer.enhance(2)
    # Save the preprocessed image
    preprocessed_image_path = image_path.replace(".png", "_preprocessed.png")
    image.save(preprocessed_image_path)
    return preprocessed_image_path

def extract_text_from_image_using_image_path(image_path):
    try:
        image = Image.open(image_path)
        languages = 'eng+spa+ara'  
        custom_config = r'--psm 1 --oem 1'

        text = pytesseract.image_to_string(image, lang=languages, config=custom_config)

        return text
    
    except Exception as e:
        print(f"error extracting text from image: {e}")
        return "No text found in image."

def extract_text_from_image(image_data):
    try:
        image = Image.open(io.BytesIO(image_data))
        languages = 'eng+spa+ara'  
        custom_config = r'--psm 1 --oem 1'
        text = pytesseract.image_to_string(image, lang=languages, config=custom_config)
        return text
    except Exception as e:
        print(f"error extracting text from image: {e}")
        return "No text found in image"

def convert_from_ppt_to_pptx(file_path_ppt, file_path_pptx):
    import win32com.client

    powerpoint = win32com.client.Dispatch("Powerpoint.Application")
    presentation = powerpoint.Presentations.open(file_path_ppt)
    presentation.saveas(file_path_pptx, 24)
    presentation.Close()
    powerpoint.Quit()

    
    
def extract_text_and_images_from_pptx(file_path, language):
        global image_description_count, last_recharge_time
        
        load_image_description_data()
        
        text_content = []  
        slide_number = 0
        media_folder = os.path.join(os.path.dirname(file_path), 'ppt', 'media')
    
        prs = Presentation(file_path)

        for i, slide in enumerate(prs.slides):        
            slide_number += 1
            text_content.append("")
            text_content.append(f"slide {slide_number}")
            print(f"Slide {slide_number}:")

            for shapes in slide.shapes:                
                if hasattr(shapes, "text_frame"):
                    for paragraph in shapes.text_frame.paragraphs:                                
                            text = paragraph.text
                            
                            text_content.append(text)                                                                 
                            
                            

            if shapes.shape_type == MSO_SHAPE_TYPE.PICTURE:  
                                if hasattr(shapes, "image"):
                                        image_data = shapes.image.blob
                                        ocr_text = extract_text_from_image(image_data)
                                        if ocr_text:
                                                text_content.append(f"Slide {slide_number}: {ocr_text}")

            for shape in slide.shapes: 
                if hasattr(shape, 'image'):
                    image            = shape.image
                    image_filename = f"image{i}_{shape.name}."
                    image_filename += image.ext if hasattr(image, 'ext') else 'jpg'

                    image_path = os.path.join(media_folder, image_filename)
                    print(f"  Image found in slide {i + 1}")
                    print(f"    Image file: {image_path}")

                    with tempfile.NamedTemporaryFile(delete= False, suffix='.' + image.ext) as temp_file:
                        temp_filename = temp_file.name
                        temp_file.write(image.blob)

                    print(f"    Image saved to temporary file: {temp_filename}")
                    
                    if image_description_count > 0:

                        describe_the_image = describe_image(temp_filename, language)

                        if describe_the_image:
                            text_content.append(f"Image description: {describe_the_image}")
                            image_description_count -= 1
                            save_image_description_data()
                            
                    else:
                        ocr_text= extract_text_from_image_using_image_path(temp_filename)
                        if ocr_text:
                            text_content.append(f"Slide {slide_number}: OCR Text from Image: {ocr_text}")
                                
                    os.remove(temp_filename)
        text_content.insert(0, f"Remaining image descriptions: {image_description_count}/20")
        
        if image_description_count == 0:
            text_content.insert(1, "You have used all your image descriptions. Please wait 4 hours to get five new image descriptions.")
    
        
        return text_content 

def extract_text_from_pdf(pdf_path):
    try:
        text = ""
        with fitz.open(pdf_path) as pdf_document:
            for page_number, page in enumerate(pdf_document, start=1):
                if page_number != 1:
                    text += "\n"
                text += f"Page {page_number}\n\n"
                
                # Extract text without using OCR
                page_text = page.get_text("text")
                text += page_text
                
        return text
    except Exception as e:
        print(f"An error occurred: {e}")
        return "An error has occurred! Please try again."
    
    
def extract_text_from_pdf_with_ocr(pdf_path, language):
    try:
        if language == "English":
            ocr_lang = "eng"
        elif language == "Arabic":
            ocr_lang = "ara"
        elif language == "Spanish":
            ocr_lang = "spa"
        
        text = ""
        with fitz.open(pdf_path) as pdf_document:
            for page_number, page in enumerate(pdf_document, start=1):
                if page_number != 1:
                    text += "\n"
                text += f"Page {page_number}\n\n"

                if page.get_pixmap():
                    image = page.get_pixmap()
                    pil_image = Image.frombytes("RGB", (image.width, image.height), image.samples)
                    page_text = pytesseract.image_to_string(pil_image, lang=ocr_lang)
                    text += page_text
                else:
                    text += page.get_text()
        return text
    except Exception as e:
        print(f"An error occurred: {e}")
        return "An error has occurred! Please try again."

